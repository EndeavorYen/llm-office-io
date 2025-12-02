#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Interactive PowerPoint Editor
強大的 PowerPoint 互動式編輯工具
支援透過自然語言指令修改 PPT 內容
"""

from typing import Optional, List
import os
import sys
import argparse

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from .constants import (
    MAX_CONTENT_PREVIEW,
    MAX_PREVIEW_LINES,
    MAX_TEXT_DISPLAY,
    SUCCESS_SYMBOL,
    ERROR_SYMBOL,
    WARNING_SYMBOL,
    PPT_EXTENSION,
    DEFAULT_LAYOUT_INDEX
)


class PPTEditor:
    """PowerPoint 編輯器類"""
    
    def __init__(self, filepath: str) -> None:
        """初始化 PowerPoint 編輯器
        
        Args:
            filepath: PowerPoint 檔案路徑
            
        Raises:
            FileNotFoundError: 當檔案不存在時
            ValueError: 當檔案格式不支援時
            RuntimeError: 當無法開啟簡報時
        """
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"檔案不存在: {filepath}")
        
        if not filepath.endswith(PPT_EXTENSION):
            raise ValueError(f"不支援的檔案格式，需要 {PPT_EXTENSION}: {filepath}")
        
        try:
            self.filepath = filepath
            self.prs = Presentation(filepath)
        except Exception as e:
            raise RuntimeError(f"無法開啟簡報: {e}") from e
    
    def save(self, output_path: Optional[str] = None) -> None:
        """儲存簡報
        
        Args:
            output_path: 輸出路徑，None 表示覆蓋原檔案
        """
        save_path = output_path or self.filepath
        try:
            self.prs.save(save_path)
            print(f"{SUCCESS_SYMBOL} 簡報已儲存: {save_path}")
        except Exception as e:
            print(f"{ERROR_SYMBOL} 儲存失敗: {e}")
            raise
    
    def list_slides(self) -> None:
        """列出所有投影片的標題和內容概要"""
        print(f"\n=== 簡報結構 (共 {len(self.prs.slides)} 張投影片) ===\n")
        
        for i, slide in enumerate(self.prs.slides, 1):
            title = self._get_slide_title(slide)
            print(f"[投影片 {i}] {title if title else '(無標題)'}")
            
            # 列出內容摘要
            content_preview = self._get_slide_content_preview(slide)
            if content_preview:
                for line in content_preview[:MAX_PREVIEW_LINES]:
                    print(f"  • {line[:MAX_CONTENT_PREVIEW]}...")
            print()
    
    def _get_slide_title(self, slide) -> Optional[str]:
        """取得投影片標題
        
        Args:
            slide: 投影片物件
            
        Returns:
            Optional[str]: 標題文字，None 表示無標題
        """
        if slide.shapes.title:
            return slide.shapes.title.text
        return None
    
    def _get_slide_content_preview(self, slide) -> List[str]:
        """取得投影片內容預覽
        
        Args:
            slide: 投影片物件
            
        Returns:
            List[str]: 內容行列表
        """
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                lines = shape.text.strip().split('\n')
                content.extend([line for line in lines if line.strip()])
        return content
    
    def replace_text(
        self, 
        old_text: str, 
        new_text: str, 
        slide_number: Optional[int] = None
    ) -> int:
        """替換文字
        
        Args:
            old_text: 要替換的文字
            new_text: 新文字
            slide_number: 指定投影片編號（從1開始），None表示全部
            
        Returns:
            int: 實際替換的次數
        """
        if not old_text:
            print(f"{ERROR_SYMBOL} 要替換的文字不能為空")
            return 0
        
        if slide_number is not None:
            if not self._validate_slide_number(slide_number):
                return 0
            slides_to_process = [self.prs.slides[slide_number - 1]]
        else:
            slides_to_process = self.prs.slides
        
        replaced_count = 0
        
        for slide in slides_to_process:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
                                replaced_count += 1
                
                # 處理表格
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if old_text in cell.text:
                                cell.text = cell.text.replace(old_text, new_text)
                                replaced_count += 1
        
        if replaced_count > 0:
            scope = f"投影片 {slide_number}" if slide_number else "所有投影片"
            print(f"{SUCCESS_SYMBOL} 在{scope}中替換了 {replaced_count} 處「{old_text}」→「{new_text}」")
        else:
            print(f"{ERROR_SYMBOL} 找不到「{old_text}」")
        
        return replaced_count
    
    def update_slide_title(self, slide_number: int, new_title: str) -> bool:
        """更新指定投影片的標題
        
        Args:
            slide_number: 投影片編號 (從1開始)
            new_title: 新標題
            
        Returns:
            bool: 是否更新成功
        """
        if not self._validate_slide_number(slide_number):
            return False
        
        slide = self.prs.slides[slide_number - 1]
        if slide.shapes.title:
            old_title = slide.shapes.title.text
            slide.shapes.title.text = new_title
            print(f"{SUCCESS_SYMBOL} 投影片 {slide_number} 標題已更新")
            print(f"  舊標題: {old_title}")
            print(f"  新標題: {new_title}")
            return True
        else:
            print(f"{ERROR_SYMBOL} 投影片 {slide_number} 沒有標題框")
            return False
    
    def add_text_to_slide(
        self, 
        slide_number: int, 
        text: str, 
        position: str = 'body'
    ) -> bool:
        """在投影片中添加文字
        
        Args:
            slide_number: 投影片編號 (從1開始)
            text: 要添加的文字
            position: 'title' 或 'body'
            
        Returns:
            bool: 是否添加成功
        """
        if not self._validate_slide_number(slide_number):
            return False
        
        slide = self.prs.slides[slide_number - 1]
        
        if position == 'title' and slide.shapes.title:
            slide.shapes.title.text = text
            print(f"{SUCCESS_SYMBOL} 已設定投影片 {slide_number} 的標題")
        else:
            # 添加到內容區
            text_added = False
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape != slide.shapes.title:
                    p = shape.text_frame.add_paragraph()
                    p.text = text
                    text_added = True
                    print(f"{SUCCESS_SYMBOL} 已在投影片 {slide_number} 添加文字")
                    break
            
            if not text_added:
                print(f"{WARNING_SYMBOL} 投影片 {slide_number} 沒有適合的文字框")
                return False
        
        return True
    
    def delete_slide(self, slide_number: int) -> bool:
        """刪除指定投影片
        
        Args:
            slide_number: 投影片編號 (從1開始)
            
        Returns:
            bool: 是否刪除成功
        """
        if not self._validate_slide_number(slide_number):
            return False
        
        title = self._get_slide_title(self.prs.slides[slide_number - 1])
        
        try:
            # 刪除投影片
            rId = self.prs.slides._sldIdLst[slide_number - 1].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[slide_number - 1]
            
            print(f"{SUCCESS_SYMBOL} 已刪除投影片 {slide_number}: {title if title else '(無標題)'}")
            return True
        except Exception as e:
            print(f"{ERROR_SYMBOL} 刪除失敗: {e}")
            return False
    
    def add_slide(self, title: str, layout_index: int = DEFAULT_LAYOUT_INDEX):
        """新增投影片
        
        Args:
            title: 投影片標題
            layout_index: 版面配置索引（預設為1，通常是標題+內容）
            
        Returns:
            Slide: 新增的投影片物件
        """
        if layout_index >= len(self.prs.slide_layouts):
            layout_index = DEFAULT_LAYOUT_INDEX
        
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        slide_num = len(self.prs.slides)
        print(f"{SUCCESS_SYMBOL} 已新增投影片 {slide_num}: {title}")
        return slide
    
    def set_font(
        self, 
        slide_number: int, 
        font_name: str, 
        font_size: Optional[int] = None
    ) -> bool:
        """設定投影片字體
        
        Args:
            slide_number: 投影片編號 (從1開始)
            font_name: 字體名稱
            font_size: 字體大小（pt）
            
        Returns:
            bool: 是否設定成功
        """
        if not self._validate_slide_number(slide_number):
            return False
        
        slide = self.prs.slides[slide_number - 1]
        changed = 0
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = font_name
                        if font_size:
                            run.font.size = Pt(font_size)
                        changed += 1
        
        size_info = f" ({font_size}pt)" if font_size else ""
        print(f"{SUCCESS_SYMBOL} 投影片 {slide_number} 已更新字體: {font_name}{size_info}")
        return True
    
    def get_slide_info(self, slide_number: int) -> None:
        """取得投影片詳細資訊
        
        Args:
            slide_number: 投影片編號 (從1開始)
        """
        if not self._validate_slide_number(slide_number):
            return
        
        slide = self.prs.slides[slide_number - 1]
        title = self._get_slide_title(slide)
        
        print(f"\n=== 投影片 {slide_number} 詳細資訊 ===")
        print(f"標題: {title if title else '(無標題)'}")
        print(f"形狀數量: {len(slide.shapes)}")
        print(f"\n內容:")
        
        for i, shape in enumerate(slide.shapes, 1):
            if hasattr(shape, "text") and shape.text:
                print(f"\n[形狀 {i}] {shape.shape_type}")
                print(shape.text[:MAX_TEXT_DISPLAY])
                if len(shape.text) > MAX_TEXT_DISPLAY:
                    print("...")
    
    def _validate_slide_number(self, slide_number: int) -> bool:
        """驗證投影片編號是否有效
        
        Args:
            slide_number: 投影片編號 (從1開始)
            
        Returns:
            bool: 是否有效
        """
        if slide_number < 1 or slide_number > len(self.prs.slides):
            print(f"{ERROR_SYMBOL} 投影片編號 {slide_number} 不存在（有效範圍: 1-{len(self.prs.slides)}）")
            return False
        return True


def main() -> None:
    """主函數"""
    parser = argparse.ArgumentParser(
        description='PowerPoint 互動式編輯器',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog='''
範例:
  # 列出所有投影片
  python ppt_editor.py presentation.pptx list
  
  # 替換文字
  python ppt_editor.py presentation.pptx replace "舊文字" "新文字"
  
  # 更新標題
  python ppt_editor.py presentation.pptx update-title 3 "新標題"
  
  # 新增投影片
  python ppt_editor.py presentation.pptx add-slide "新投影片標題"
  
  # 刪除投影片
  python ppt_editor.py presentation.pptx delete-slide 5
        '''
    )
    
    parser.add_argument('file', help='PowerPoint 檔案路徑')
    parser.add_argument('--output', '-o', help='輸出檔案路徑（不指定則覆蓋原檔案）')
    
    subparsers = parser.add_subparsers(dest='command', help='編輯命令')
    
    # list: 列出投影片
    subparsers.add_parser('list', help='列出所有投影片')
    
    # replace: 替換文字
    replace_parser = subparsers.add_parser('replace', help='替換文字')
    replace_parser.add_argument('old', help='要替換的文字')
    replace_parser.add_argument('new', help='新文字')
    replace_parser.add_argument('--slide', type=int, help='指定投影片編號（不指定則全部）')
    
    # update-title: 更新標題
    title_parser = subparsers.add_parser('update-title', help='更新投影片標題')
    title_parser.add_argument('slide', type=int, help='投影片編號')
    title_parser.add_argument('title', help='新標題')
    
    # add-text: 添加文字
    addtext_parser = subparsers.add_parser('add-text', help='在投影片中添加文字')
    addtext_parser.add_argument('slide', type=int, help='投影片編號')
    addtext_parser.add_argument('text', help='要添加的文字')
    
    # add-slide: 新增投影片
    addslide_parser = subparsers.add_parser('add-slide', help='新增投影片')
    addslide_parser.add_argument('title', help='投影片標題')
    addslide_parser.add_argument('--layout', type=int, default=DEFAULT_LAYOUT_INDEX, 
                                help='版面配置索引（預設1）')
    
    # delete-slide: 刪除投影片
    delete_parser = subparsers.add_parser('delete-slide', help='刪除投影片')
    delete_parser.add_argument('slide', type=int, help='投影片編號')
    
    # info: 查看投影片資訊
    info_parser = subparsers.add_parser('info', help='查看投影片詳細資訊')
    info_parser.add_argument('slide', type=int, help='投影片編號')
    
    # set-font: 設定字體
    font_parser = subparsers.add_parser('set-font', help='設定投影片字體')
    font_parser.add_argument('slide', type=int, help='投影片編號')
    font_parser.add_argument('font', help='字體名稱')
    font_parser.add_argument('--size', type=int, help='字體大小（pt）')
    
    args = parser.parse_args()
    
    if not args.command:
        parser.print_help()
        return
    
    # 載入簡報
    try:
        editor = PPTEditor(args.file)
    except (FileNotFoundError, ValueError, RuntimeError) as e:
        print(f"{ERROR_SYMBOL} {e}")
        sys.exit(1)
    
    # 執行命令
    try:
        if args.command == 'list':
            editor.list_slides()
            return
        
        elif args.command == 'replace':
            editor.replace_text(args.old, args.new, args.slide)
        
        elif args.command == 'update-title':
            editor.update_slide_title(args.slide, args.title)
        
        elif args.command == 'add-text':
            editor.add_text_to_slide(args.slide, args.text)
        
        elif args.command == 'add-slide':
            editor.add_slide(args.title, args.layout)
        
        elif args.command == 'delete-slide':
            editor.delete_slide(args.slide)
        
        elif args.command == 'info':
            editor.get_slide_info(args.slide)
            return
        
        elif args.command == 'set-font':
            editor.set_font(args.slide, args.font, args.size)
        
        # 儲存
        editor.save(args.output)
        
    except Exception as e:
        print(f"{ERROR_SYMBOL} 操作失敗: {e}")
        sys.exit(1)


if __name__ == '__main__':
    main()
